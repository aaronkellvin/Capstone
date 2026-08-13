import io
from pathlib import Path

MIN_CHARS = 200
ALLOWED_EXT = {".pdf", ".docx", ".pptx", ".txt"}


class ExtractError(Exception):
    pass


def extract_text(filename: str, data: bytes) -> str:
    ext = Path(filename).suffix.lower()
    if ext not in ALLOWED_EXT:
        raise ExtractError("Please upload a PDF, DOCX, PPTX, or TXT file.")

    if ext == ".txt":
        text = data.decode("utf-8", errors="ignore")
    elif ext == ".pdf":
        text = _pdf(data)
    elif ext == ".docx":
        text = _docx(data)
    else:
        text = _pptx(data)

    cleaned = " ".join(text.split())
    if len(cleaned) < MIN_CHARS:
        raise ExtractError(
            "This file does not have enough extractable text. "
            "Scanned or image-only PDFs are not supported. Upload a text-based file."
        )
    return cleaned[:50000]


def _pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise ExtractError("PDF support is not installed. Run: pip install pypdf") from exc
    reader = PdfReader(io.BytesIO(data))
    if len(reader.pages) > 100:
        raise ExtractError("PDFs are limited to 100 pages for this pilot.")
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(data: bytes) -> str:
    try:
        from docx import Document
    except ImportError as exc:
        raise ExtractError("DOCX support is not installed. Run: pip install python-docx") from exc
    document = Document(io.BytesIO(data))
    return "\n".join(p.text for p in document.paragraphs)


def _pptx(data: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ExtractError("PPTX support is not installed. Run: pip install python-pptx") from exc
    pres = Presentation(io.BytesIO(data))
    if len(pres.slides) > 100:
        raise ExtractError("Presentations are limited to 100 slides for this pilot.")
    chunks = []
    for slide in pres.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                chunks.append(shape.text)
    return "\n".join(chunks)
